"""Slide builder — converts video script JSON into PNG slides via python-pptx + LibreOffice.

Pipeline:
  slide data dict → python-pptx PPTX file → LibreOffice CLI → individual PNG files
"""

import os
import subprocess
import tempfile
from pathlib import Path

import structlog
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt

logger = structlog.get_logger(__name__)

# Slide dimensions (16:9 widescreen)
SLIDE_WIDTH = Inches(13.33)
SLIDE_HEIGHT = Inches(7.5)

# Color scheme — dark editorial theme
BG_COLOR = RGBColor(0x0A, 0x0A, 0x0F)
ACCENT_COLOR = RGBColor(0x3B, 0x82, 0xF6)   # Electric blue
TEXT_COLOR = RGBColor(0xFF, 0xFF, 0xFF)
SUBTITLE_COLOR = RGBColor(0x94, 0xA3, 0xB8)
CODE_BG = RGBColor(0x1E, 0x29, 0x3B)


def build_slides(script: dict, output_dir: str) -> list[str]:
    """Build PNG slides from the video script.

    Args:
        script: Parsed video script dict with 'slides' list.
        output_dir: Directory to write PNG files.

    Returns:
        Ordered list of PNG file paths.
    """
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    blank_layout = prs.slide_layouts[6]  # Completely blank layout

    for slide_data in script.get("slides", []):
        slide = prs.slides.add_slide(blank_layout)
        _fill_slide_background(slide, BG_COLOR)

        slide_type = slide_data.get("slide_type", "concept")
        if slide_type == "title":
            _build_title_slide(slide, slide_data)
        elif slide_type == "code":
            _build_code_slide(slide, slide_data)
        else:
            _build_content_slide(slide, slide_data)

    # Save PPTX to temp file
    pptx_path = os.path.join(output_dir, "slides.pptx")
    prs.save(pptx_path)
    logger.info("PPTX built", slides=len(script.get("slides", [])), path=pptx_path)

    # Export slides to PNG via LibreOffice
    png_paths = _export_pptx_to_pngs(pptx_path, output_dir)
    return png_paths


def _fill_slide_background(slide, color: RGBColor) -> None:
    """Fill slide background with a solid color."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_text_box(slide, text: str, left, top, width, height, size: int,
                  bold: bool = False, color: RGBColor = TEXT_COLOR,
                  wrap: bool = True) -> None:
    """Add a text box to a slide."""
    from pptx.util import Emu
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.text = text
    run = p.runs[0] if p.runs else p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def _build_title_slide(slide, data: dict) -> None:
    """Build a title card slide."""
    _add_text_box(
        slide, data.get("title", ""), Inches(1), Inches(2.5),
        Inches(11.33), Inches(1.5), size=48, bold=True, color=TEXT_COLOR
    )
    bullets = data.get("bullet_points", [])
    if bullets:
        subtitle = bullets[0] if bullets else ""
        _add_text_box(
            slide, subtitle, Inches(1), Inches(4.3),
            Inches(11.33), Inches(1), size=24, color=SUBTITLE_COLOR
        )
    # Accent line
    from pptx.util import Emu
    line = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(1), Inches(4.1), Inches(2), Pt(3)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT_COLOR
    line.line.fill.background()


def _build_content_slide(slide, data: dict) -> None:
    """Build a standard concept/example/summary slide."""
    # Title
    _add_text_box(
        slide, data.get("title", ""), Inches(0.5), Inches(0.3),
        Inches(12.33), Inches(1), size=32, bold=True, color=ACCENT_COLOR
    )
    # Bullet points
    bullets = data.get("bullet_points", [])
    if bullets:
        bullet_text = "\n".join(f"• {b}" for b in bullets)
        _add_text_box(
            slide, bullet_text, Inches(0.5), Inches(1.5),
            Inches(12.33), Inches(5.5), size=20, color=TEXT_COLOR
        )


def _build_code_slide(slide, data: dict) -> None:
    """Build a slide with a code snippet."""
    _add_text_box(
        slide, data.get("title", ""), Inches(0.5), Inches(0.3),
        Inches(12.33), Inches(0.8), size=28, bold=True, color=ACCENT_COLOR
    )
    code = data.get("code_snippet", "") or ""
    if code:
        # Code box with dark background
        code_box = slide.shapes.add_shape(1, Inches(0.5), Inches(1.2), Inches(12.33), Inches(5.8))
        code_box.fill.solid()
        code_box.fill.fore_color.rgb = CODE_BG
        code_box.line.fill.background()

        from pptx.util import Emu
        txBox = slide.shapes.add_textbox(Inches(0.7), Inches(1.4), Inches(12), Inches(5.4))
        tf = txBox.text_frame
        tf.word_wrap = False
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = code
        run.font.size = Pt(14)
        run.font.color.rgb = RGBColor(0xA7, 0xF3, 0xD0)  # mint green for code
        run.font.name = "Courier New"


def _export_pptx_to_pngs(pptx_path: str, output_dir: str) -> list[str]:
    """Use LibreOffice CLI to convert PPTX to individual PNG files.

    Args:
        pptx_path: Path to the PPTX file.
        output_dir: Directory to write PNG files.

    Returns:
        Sorted list of PNG file paths.
    """
    try:
        result = subprocess.run(
            [
                "libreoffice",
                "--headless",
                "--convert-to", "png",
                "--outdir", output_dir,
                pptx_path,
            ],
            capture_output=True,
            text=True,
            timeout=120,
        )
        if result.returncode != 0:
            raise RuntimeError(
                f"LibreOffice conversion failed: {result.stderr}"
            )
    except FileNotFoundError:
        # LibreOffice not installed — generate placeholder PNGs with Pillow
        logger.warning("LibreOffice not found — generating placeholder slides with Pillow")
        return _generate_placeholder_pngs(pptx_path, output_dir)

    # LibreOffice outputs files named like "slides.png" (single) or "slides-slide1.png"
    pngs = sorted(Path(output_dir).glob("slides*.png"))
    if not pngs:
        raise RuntimeError("No PNG files found after LibreOffice conversion")

    logger.info("PPTX converted to PNGs", count=len(pngs))
    return [str(p) for p in pngs]


def _generate_placeholder_pngs(pptx_path: str, output_dir: str) -> list[str]:
    """Fallback: generate dark-background PNG slides with Pillow."""
    from PIL import Image, ImageDraw, ImageFont
    from pptx import Presentation as _Prs

    prs = _Prs(pptx_path)
    paths = []

    for i, slide in enumerate(prs.slides):
        img = Image.new("RGB", (1920, 1080), (10, 10, 15))
        draw = ImageDraw.Draw(img)

        # Draw blue top bar
        draw.rectangle([0, 0, 1920, 8], fill=(59, 130, 246))

        # Extract text from slide shapes
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    t = para.text.strip()
                    if t:
                        texts.append(t)

        y = 80
        for j, text in enumerate(texts[:20]):
            font_size = 48 if j == 0 else 28
            try:
                draw.text((80, y), text[:100], fill=(255, 255, 255))
            except Exception:
                pass
            y += font_size + 16

        png_path = os.path.join(output_dir, f"slide_{i+1:03d}.png")
        img.save(png_path, "PNG")
        paths.append(png_path)

    return sorted(paths)
